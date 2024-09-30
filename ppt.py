# Re-run the code to create the PowerPoint presentation with the same content
import pptx
from pptx.util import Inches

# Create a presentation object
prs = pptx.Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]  # Title Slide layout
bullet_slide_layout = prs.slide_layouts[1]  # Title and Content layout

# Add title slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Integrating Emerging Technologies in Resource-Constrained Educational Institutions"
subtitle.text = "An Evaluation of Augmented Reality, Virtual Reality, and Artificial Intelligence\nResearch Proposal\nYour Name\nDate"

# Slide 2: Introduction
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Introduction"
content.text = ("Objective: To evaluate the effectiveness and scalability of AR, VR, and AI technologies in resource-constrained educational settings.\n"
                "Importance: Addressing educational disparities and promoting equity through innovative technology.")

# Slide 3: Research Question
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Research Question"
content.text = ("Main Question: How can emerging technologies like AR, VR, and AI be developed further to operate effectively under limited resources in any educational institution, ensuring diverse accessibility and minimizing technological disparities?")

# Slide 4: Literature Review
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Literature Review"
content.text = ("Overview: Summary of key studies on AR, VR, and AI in education.\n"
                "Gaps: Limited research on implementation in resource-constrained settings and ensuring diverse accessibility.")

# Slide 5: Problem Definition
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Problem Definition"
content.text = ("Gap: Lack of empirical research on the practical impacts of AR, VR, and AI in under-resourced schools.\n"
                "Need: Developing scalable, cost-effective solutions to enhance educational equity.")

# Slide 6: Research Method
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Research Method"
content.text = ("Design: Mixed-methods approach combining qualitative and quantitative methods.\n"
                "Data Collection: Surveys, interviews, experimental data, and observation logs.\n"
                "Sample: Stratified random sampling across different educational levels and locations.")

# Slide 7: Data Sources
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Data Sources"
content.text = ("Surveys: Distributed to educators, students, and administrators.\n"
                "Interviews: In-depth discussions with a representative sample.\n"
                "Experimental Data: Performance metrics from controlled experiments.\n"
                "Observation Logs: Detailed interaction logs during implementation.")

# Slide 8: Data Collection Methods
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Data Collection Methods"
content.text = ("Surveys: Online surveys capturing quantitative data.\n"
                "Interviews: Semi-structured interviews for qualitative insights.\n"
                "Experimental Data Collection: Metrics on test scores, task completion rates, and attendance.\n"
                "Observation Logs: Objective measures of user engagement and behavior.")

# Slide 9: Reliability and Legitimacy of Data
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Reliability and Legitimacy of Data"
content.text = ("Informed Consent: Ensuring voluntary participation and confidentiality.\n"
                "Authorization: Compliance with data privacy regulations.\n"
                "Ethical Standards: Adherence to IRB guidelines and ethical considerations.\n"
                "Data Integrity: Secure storage and quality control checks.")

# Slide 10: Ability to Meet Research Needs
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Ability to Meet Research Needs"
content.text = ("Comprehensive Approach: Combining surveys, interviews, experiments, and logs for a holistic view.\n"
                "Objective Measures: Direct assessment of AR, VR, and AI impacts on engagement and outcomes.\n"
                "Robustness: Triangulation of findings for increased validity.")

# Slide 11: Data Analysis Methods
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Data Analysis Methods"
content.text = ("Quantitative Analysis: Descriptive and inferential statistics, regression analysis.\n"
                "Qualitative Analysis: Thematic analysis and systematic coding.\n"
                "Triangulation: Integrating quantitative and qualitative findings for comprehensive insights.")

# Slide 12: Discussion
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Discussion"
content.text = ("Potential Impact: Enhancing educational equity, technological innovation, and policy guidance.\n"
                "Positive Impact: Improved engagement and outcomes, informed policies, and innovative practices.\n"
                "No Impact: Identification of alternative strategies and focus on foundational resource needs.")

# Slide 13: Contribution to the Field
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Contribution to the Field"
content.text = ("Framework for Integration: Practical considerations for cost, scalability, and usability.\n"
                "Best Practices: Guidelines for educators and policymakers.\n"
                "Enhanced Understanding: Leveraging technologies to promote educational equity.\n"
                "Policy Recommendations: Supporting widespread adoption in under-resourced schools.")

# Slide 14: Conclusion
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusion"
content.text = ("Summary: Addressing gaps in research on AR, VR, and AI in resource-constrained settings.\n"
                "Future Directions: Long-term studies, scalability, and user-centered design.\n"
                "Significance: Contributing to a more equitable and effective educational environment.")

# Slide 15: Q&A
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Q&A"
content.text = "Questions: Open the floor for questions and discussion."

# Save the presentation
presentation_file = "./Research_Proposal_Presentation.pptx"
prs.save(presentation_file)

presentation_file
